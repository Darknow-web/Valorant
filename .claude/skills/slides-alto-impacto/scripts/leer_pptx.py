#!/usr/bin/env python3
"""Lee un .pptx existente y lo convierte al mismo formato de shapes que usa
render_html.cjs, para poder VERLO y auditarlo.

    python3 leer_pptx.py presentacion.pptx -o salida/original

Aquí LibreOffice está incompleto y no convierte .pptx, así que este es el único
modo de mirar una presentación ajena: se reconstruye desde su estructura
(posiciones, textos, colores, formas, imágenes) y se dibuja con el renderizador
que ya existe.

Produce:
  salida/original/deck.plan.json   compatible con render_html.cjs y auditar.py
  salida/original/img/*            las imágenes embebidas
"""
import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

EMU_POR_PULGADA = 914400


def pulgadas(v):
    return round((v or 0) / EMU_POR_PULGADA, 4)


def color_de(objeto_color, defecto=None):
    """Hex de un color de python-pptx, o el defecto si es de tema o no hay."""
    try:
        if objeto_color is None:
            return defecto
        if objeto_color.type is None:
            return defecto
        # Los colores de tema no exponen rgb; se resuelven fuera.
        return str(objeto_color.rgb)
    except Exception:
        return defecto


def relleno_de(forma):
    try:
        f = forma.fill
        if f.type is None or f.type == 5:      # sin relleno / hereda
            return None
        if f.type == 1:                        # sólido
            return color_de(f.fore_color)
    except Exception:
        pass
    return None


def linea_de(forma):
    try:
        ln = forma.line
        if ln.fill.type in (None, 5):
            return None
        color = color_de(ln.color)
        if not color:
            return None
        return {"color": color, "width": round((ln.width or 0) / 12700, 2) or 1}
    except Exception:
        return None


def parrafos_de(marco):
    """Texto por párrafo, con el formato dominante de cada uno."""
    salida = []
    for p in marco.paragraphs:
        runs = [r for r in p.runs if r.text]
        texto = "".join(r.text for r in p.runs)
        if not texto.strip():
            continue
        r0 = runs[0] if runs else None
        f = r0.font if r0 is not None else None
        salida.append({
            "texto": texto,
            "pt": (f.size.pt if f is not None and f.size else None),
            "bold": bool(f.bold) if f is not None else False,
            "italic": bool(f.italic) if f is not None else False,
            "color": color_de(f.color, None) if f is not None else None,
            "nivel": p.level,
            "align": str(p.alignment).split(".")[0].lower() if p.alignment else None,
        })
    return salida


def leer(ruta, dir_salida):
    pres = Presentation(str(ruta))
    ancho = pulgadas(pres.slide_width)
    alto = pulgadas(pres.slide_height)
    dir_img = dir_salida / "img"
    dir_img.mkdir(parents=True, exist_ok=True)

    slides = []
    n_img = 0
    for i, slide in enumerate(pres.slides, 1):
        shapes = []
        fondo = None
        try:
            if slide.background.fill.type == 1:
                fondo = color_de(slide.background.fill.fore_color)
        except Exception:
            pass
        shapes.append({"tipo": "background", "opciones": {"color": fondo or "FFFFFF"}})

        for j, forma in enumerate(slide.shapes):
            nombre = f"s{i}_orig{j}"
            o = {
                "x": pulgadas(forma.left), "y": pulgadas(forma.top),
                "w": pulgadas(forma.width), "h": pulgadas(forma.height),
                "objectName": nombre,
            }

            if forma.shape_type == 13 or forma.__class__.__name__ == "Picture":
                try:
                    img = forma.image
                    destino = dir_img / f"img{n_img:03d}.{img.ext}"
                    destino.write_bytes(img.blob)
                    n_img += 1
                    shapes.append({"tipo": "addImage",
                                   "opciones": {**o, "path": str(destino)}})
                except Exception:
                    shapes.append({"tipo": "addShape", "arg0": "rect",
                                   "opciones": {**o, "fill": {"color": "CCCCCC"}}})
                continue

            if forma.has_table:
                tabla = forma.table
                datos = [[{"text": c.text, "options": {}} for c in fila.cells]
                         for fila in tabla.rows]
                shapes.append({"tipo": "addTable", "datos": datos,
                               "opciones": {**o, "rowH": round(
                                   (o["h"] or 1) / max(len(datos), 1), 3)}})
                continue

            if forma.has_chart:
                shapes.append({"tipo": "addShape", "arg0": "rect",
                               "opciones": {**o, "fill": {"color": "DDDDDD"},
                                            "__grafico": True}})
                continue

            relleno = relleno_de(forma)
            linea = linea_de(forma)
            if relleno or linea:
                op = {**o, "objectName": nombre + "_fondo"}
                if relleno:
                    op["fill"] = {"color": relleno}
                if linea:
                    op["line"] = linea
                shapes.append({"tipo": "addShape", "arg0": "rect", "opciones": op})

            if forma.has_text_frame:
                for k, p in enumerate(parrafos_de(forma.text_frame)):
                    # Cada párrafo va como su propio cuadro, apilado dentro del
                    # marco original: es lo que permite medirlos por separado.
                    n_par = max(len(parrafos_de(forma.text_frame)), 1)
                    alto_p = (o["h"] or 0.4) / n_par
                    shapes.append({
                        "tipo": "addText",
                        "arg0": p["texto"],
                        "opciones": {
                            **o,
                            "y": round((o["y"] or 0) + alto_p * k, 4),
                            "h": round(alto_p, 4),
                            "objectName": f"{nombre}_p{k}",
                            "fontSize": p["pt"] or 18,
                            "bold": p["bold"], "italic": p["italic"],
                            "color": p["color"] or "222222",
                            # Sin color explícito lo hereda del tema del .pptx,
                            # que no se puede resolver desde el archivo.
                            "__color_heredado": p["color"] is None,
                            "align": p["align"] if p["align"] in
                                     ("center", "right") else "left",
                        },
                    })

        notas = ""
        try:
            if slide.has_notes_slide:
                notas = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass

        slides.append({
            "indice": i, "arquetipo": "importado", "objetos": [],
            "transicion": "ninguna", "animacion": {},
            "dato": None, "grafico": None, "items": None,
            "notas": notas, "shapes": shapes,
        })

    # Colores dominantes, para que el renderizador tenga con qué pintar.
    textos = [sh["opciones"].get("color") for s in slides for sh in s["shapes"]
              if sh["tipo"] == "addText"]
    fondos = [sh["opciones"].get("color") for s in slides for sh in s["shapes"]
              if sh["tipo"] == "background"]
    dominante = max(set(fondos), key=fondos.count) if fondos else "FFFFFF"

    return {
        "titulo": ruta.stem,
        "tema": "importado",
        "origen": str(ruta),
        "lienzo": {"ancho_in": ancho, "alto_in": alto},
        "colores": {
            "modo": "claro", "fondo": dominante, "fondo_alt": dominante,
            "superficie": "FFFFFF", "borde": "CCCCCC",
            "texto": (max(set(textos), key=textos.count) if textos else "222222"),
            "texto_suave": "666666", "acento": "2563EB", "acento2": "7C3AED",
            "serie": ["2563EB", "7C3AED", "059669", "D97706", "DC2626", "0891B2"],
            "fuente_titulo": "Inter", "fuente_cuerpo": "Inter",
        },
        "slides": slides,
    }


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("pptx")
    ap.add_argument("-o", "--salida", default="original")
    args = ap.parse_args()

    ruta = Path(args.pptx)
    if not ruta.exists():
        sys.exit(f"No existe: {ruta}")
    salida = Path(args.salida)
    salida.mkdir(parents=True, exist_ok=True)

    plan = leer(ruta, salida)
    destino = salida / "deck.plan.json"
    destino.write_text(json.dumps(plan, ensure_ascii=False, indent=2),
                       encoding="utf-8")

    n_shapes = sum(len(s["shapes"]) for s in plan["slides"])
    print(f"  {len(plan['slides'])} slides · {n_shapes} elementos")
    print(f"  lienzo {plan['lienzo']['ancho_in']}x{plan['lienzo']['alto_in']} in")
    print(f"  {destino}")
    print("\n  Para verlo:")
    print(f"    node scripts/render_html.cjs {destino} -o {salida}/vista.html")
    print(f"    node scripts/capturar.cjs {salida}/vista.html -o {salida}/png --grid")


if __name__ == "__main__":
    main()
